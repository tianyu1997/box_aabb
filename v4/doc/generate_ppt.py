"""
SafeBoxForest presentation generator (v10).
- 11 slides
- clear 4-minute narrative
- embedded figures + timed speaker notes
"""

from __future__ import annotations

import os
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

FONT = "Calibri"
W_IN, H_IN = 13.333, 7.5

C_BG = RGBColor(0xF8, 0xFA, 0xFD)
C_DARK = RGBColor(0x16, 0x24, 0x3D)
C_BLUE = RGBColor(0x2E, 0x5E, 0xA6)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
C_RED = RGBColor(0xC6, 0x28, 0x28)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_LGRAY = RGBColor(0xE9, 0xEE, 0xF7)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x22, 0x22, 0x22)


def _rect(slide, left, top, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _round_rect(slide, left, top, w, h, color, text, sz=16):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(sz)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    return shp


def _txt(slide, left, top, w, h, text, sz=18, color=C_BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    return tb


def _bullets(slide, left, top, w, h, items: Iterable[str], sz=18, color=C_BLACK, spacing=6):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return tb


def _header(slide, title: str):
    _rect(slide, Inches(0), Inches(0), Inches(W_IN), Inches(1.05), C_DARK)
    _txt(slide, Inches(0.7), Inches(0.2), Inches(12), Inches(0.6), title, sz=32, color=C_WHITE, bold=True)


def _set_notes(slides, notes):
    for idx, slide in enumerate(slides):
        if idx >= len(notes):
            break
        ns = slide.notes_slide
        ns.notes_text_frame.text = notes[idx]


def _img(slide, file_name: str, left, top, width=None, height=None):
    img_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), file_name)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)


def build(notes_lang: str = "zh") -> str:
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    BLANK = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_DARK)
    _rect(s, Inches(1.0), Inches(2.2), Inches(2.4), Pt(4), C_ORANGE)
    _txt(s, Inches(1.0), Inches(2.45), Inches(10), Inches(1.2), "SafeBoxForest", sz=56, color=C_WHITE, bold=True)
    _txt(s, Inches(1.0), Inches(3.7), Inches(11), Inches(1.1),
         "Interval-AABB Certified Motion Planning\nwith Persistent Free-Space Forest",
         sz=24, color=RGBColor(0xD6, 0xDE, 0xF0))
    _txt(s, Inches(1.0), Inches(5.45), Inches(4), Inches(0.5), "IROS 2026", sz=20, color=C_ORANGE, bold=True)

    # 2 Problem
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_WHITE)
    _header(s, "问题与切入点")
    _txt(s, Inches(0.8), Inches(1.45), Inches(5.8), Inches(0.5), "现有方法痛点", sz=23, color=C_RED, bold=True)
    _bullets(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(3.2), [
        "RRT/PRM 跨查询不可复用",
        "碰撞检测是逐点成本，缺体积证书",
        "障碍变化后常需全量重建",
    ], sz=19)
    _txt(s, Inches(7.0), Inches(1.45), Inches(5.3), Inches(0.5), "SBF 核心思想", sz=23, color=C_GREEN, bold=True)
    _bullets(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.4), [
        "用 Safe Box 覆盖关节自由空间",
        "每个 Box 带区间 FK 安全证书",
        "Forest 持久化并支持增量修复",
    ], sz=19)
    _round_rect(s, Inches(0.8), Inches(5.75), Inches(11.8), Inches(0.9), C_BLUE,
                "一句话：从“点采样搜索”升级为“可复用体积安全地图”", sz=20)

    # 3 Architecture
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_BG)
    _header(s, "三层架构")
    layers = [
        ("规划层", "Graph Search + Smoothing + Query", C_ORANGE),
        ("森林层", "SafeBoxForest + Adjacency + Incremental Update", C_BLUE),
        ("证书层", "Interval FK + AABB Collision Certificate", C_GREEN),
    ]
    for i, (t, d, c) in enumerate(layers):
        y = Inches(1.55 + i * 1.8)
        _round_rect(s, Inches(2.3), y, Inches(8.8), Inches(0.9), c, t, sz=22)
        _txt(s, Inches(2.6), y + Inches(0.95), Inches(8.2), Inches(0.45), d, sz=14, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 2:
            _txt(s, Inches(6.35), y + Inches(1.35), Inches(0.6), Inches(0.3), "▼", sz=22, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 4 Pipeline
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_WHITE)
    _header(s, "算法流程（端到端）")
    cards = [
        ("① 区间证书", "Interval FK 判定 Safe Box", C_GREEN),
        ("② 森林生长", "KD-Tree + Promote", C_BLUE),
        ("③ 图连通", "Adjacency + Bridge", C_ORANGE),
        ("④ 查询优化", "Geodesic Dijkstra + Smooth", C_DARK),
    ]
    for i, (t, d, c) in enumerate(cards):
        x = Inches(0.6 + i * 3.17)
        _round_rect(s, x, Inches(1.6), Inches(2.9), Inches(0.8), c, t, sz=16)
        _txt(s, x + Inches(0.1), Inches(2.55), Inches(2.7), Inches(1.0), d, sz=15, align=PP_ALIGN.CENTER)
    _rect(s, Inches(0.8), Inches(4.2), Inches(11.8), Inches(2.4), C_LGRAY)
    _txt(s, Inches(1.1), Inches(4.35), Inches(11.2), Inches(0.5), "关键优化", sz=22, color=C_BLUE, bold=True)
    _bullets(s, Inches(1.1), Inches(4.95), Inches(11.2), Inches(1.5), [
        "周期空间一致性：geodesic cost，避免 ±π 边界失真",
        "fallback waypoint 使用 wrap-aware nearest point",
        "多查询复用 + 增量重建提升整体吞吐",
    ], sz=17)

    # 5 Interval FK + figure
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_WHITE)
    _header(s, "Safe Box 如何判定")
    _txt(s, Inches(0.8), Inches(1.45), Inches(5.9), Inches(0.5), "区间 FK 安全证书", sz=23, color=C_GREEN, bold=True)
    _bullets(s, Inches(0.8), Inches(2.05), Inches(5.9), Inches(3.6), [
        "把关节值推广为区间并传播 DH 链",
        "得到每个 link 的保守 AABB",
        "若所有 AABB 与障碍不相交 ⇒ 整个盒子安全",
        "特性：无假阴性（安全优先）",
    ], sz=17)
    _img(s, "img_interval_fk_aabb.png", Inches(6.7), Inches(1.9), width=Inches(6.0))
    _round_rect(s, Inches(0.8), Inches(6.1), Inches(5.9), Inches(0.75), C_GREEN,
                "证书化判定：快于密集采样且可复用", sz=16)

    # 6 KD-Tree + figure
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_BG)
    _header(s, "KD-Tree 生长与缓存")
    _bullets(s, Inches(0.8), Inches(1.55), Inches(5.8), Inches(4.2), [
        "递归二分高维空间，非重叠性天然成立",
        "Find-Free-Box + Promote 找最大可行盒",
        "AABB 缓存复用显著降低 FK 计算成本",
        "持久化 hcache 支持重启后快速加载",
    ], sz=17)
    _img(s, "img_forest_final.png", Inches(6.8), Inches(1.8), width=Inches(5.8))
    _txt(s, Inches(6.8), Inches(5.95), Inches(5.8), Inches(0.35),
         "2DOF forest growth snapshot", sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 7 Incremental + figure
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_WHITE)
    _header(s, "复用与增量更新")
    _round_rect(s, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.65), C_RED, "1) Invalidate", sz=14)
    _txt(s, Inches(0.9), Inches(2.2), Inches(3.3), Inches(0.9), "仅标记受新障碍影响的盒子", sz=14)
    _round_rect(s, Inches(4.5), Inches(1.5), Inches(3.5), Inches(0.65), C_ORANGE, "2) Remove", sz=14)
    _txt(s, Inches(4.6), Inches(2.2), Inches(3.3), Inches(0.9), "删除失效盒，保留其余结构", sz=14)
    _round_rect(s, Inches(8.2), Inches(1.5), Inches(3.5), Inches(0.65), C_GREEN, "3) Regrow", sz=14)
    _txt(s, Inches(8.3), Inches(2.2), Inches(3.3), Inches(0.9), "只在空洞区域局部补种", sz=14)
    _img(s, "img_incremental_update.png", Inches(0.8), Inches(3.25), width=Inches(11.8))
    _round_rect(s, Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.6), C_BLUE,
                "结论：环境变化不再全量重建，时延更稳定", sz=15)

    # 8 Results
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_WHITE)
    _header(s, "结果总览")
    metrics = [
        ("100%", "窄通道成功率", C_GREEN),
        ("50–80×", "多查询摊销加速", C_BLUE),
        ("<2.2s", "障碍变化适应", C_ORANGE),
        ("稳定", "周期空间路径一致", C_DARK),
    ]
    for i, (v, t, c) in enumerate(metrics):
        x = Inches(0.55 + i * 3.18)
        _round_rect(s, x, Inches(1.55), Inches(2.9), Inches(1.2), c, v, sz=36)
        _txt(s, x, Inches(2.85), Inches(2.9), Inches(0.45), t, sz=15, align=PP_ALIGN.CENTER)
    _rect(s, Inches(0.8), Inches(4.0), Inches(11.8), Inches(2.75), C_LGRAY)
    _bullets(s, Inches(1.0), Inches(4.25), Inches(11.4), Inches(2.2), [
        "SBF 在复用查询和动态障碍场景更具优势",
        "RRT 起步快但跨查询不可复用，整体吞吐受限",
        "路径偏置问题已通过 geodesic + wrap-aware 修复链路处理",
    ], sz=18)

    # 9 Positioning
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_BG)
    _header(s, "方法定位")
    _txt(s, Inches(0.9), Inches(1.5), Inches(11.8), Inches(0.6),
         "SBF 的差异化：安全证书 + 可复用 + 增量更新 + 并行扩展", sz=21, color=C_BLUE, bold=True)
    _bullets(s, Inches(1.0), Inches(2.2), Inches(11.5), Inches(3.8), [
        "IRIS/GCS：高质量凸优化，但增量与工程成本较高",
        "RRT 家族：快速起步，但复用性与稳定性弱",
        "SBF：在动态、重复查询场景提供更优综合效率",
    ], sz=19)

    # 10 Conclusion
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_DARK)
    _txt(s, Inches(1.0), Inches(1.0), Inches(11), Inches(0.8), "总结", sz=42, color=C_WHITE, bold=True)
    _bullets(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.6), [
        "区间 FK 证书化判定：安全、可复用",
        "SafeBoxForest：持久化 free-space 地图",
        "动态障碍：局部增量修复替代全量重建",
        "路径质量：geodesic + wrap-aware 链路保证周期一致性",
    ], sz=22, color=C_WHITE)
    _txt(s, Inches(1.0), Inches(6.2), Inches(11), Inches(0.5),
         "Future: tighter envelope, denser scenes, TAMP integration",
         sz=17, color=RGBColor(0xD6, 0xDE, 0xF0))

    # 11 Q&A
    s = prs.slides.add_slide(BLANK)
    _rect(s, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN), C_DARK)
    _txt(s, Inches(0), Inches(2.45), Inches(W_IN), Inches(1.2), "Thank You", sz=58, color=C_WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    _txt(s, Inches(0), Inches(4.0), Inches(W_IN), Inches(0.7), "Questions?", sz=30,
         color=RGBColor(0xD6, 0xDE, 0xF0), align=PP_ALIGN.CENTER)

    if notes_lang.lower() == "en":
        notes = [
            "[20s] Core message: SBF upgrades one-shot sampling into a reusable volumetric safety map. Trigger: contribution? -> certificate + reuse + incremental update.",
            "[20s] Pain points: non-reusable queries, point-wise collision cost, full rebuild under scene changes. Trigger: why not plain RRT? -> poor cross-query throughput.",
            "[20s] Three-layer decoupling: certificate layer, forest layer, planning layer. Trigger: engineering scalability? -> each layer is independently replaceable/parallelizable.",
            "[25s] Pipeline in 4 steps, plus key fix: geodesic cost and wrap-aware waypointing. Trigger: path bias? -> fixed ±pi boundary distortion.",
            "[20s] Interval FK gives conservative but safe certificates. Trigger: over-conservative? -> possible false positives, but no missed collision.",
            "[20s] KD-tree + promote + cache reuse is the efficiency core. Trigger: why fast? -> reusable subproblems and reduced FK recomputation.",
            "[20s] Dynamic scenes use invalidate-remove-regrow locally. Trigger: dynamic capability? -> avoids full rebuild.",
            "[20s] Report only success, speed, adaptability. Trigger: path quality? -> geodesic edges + wrap-aware waypoints.",
            "[20s] Positioning is complementary: SBF vs IRIS/GCS/RRT by scenario fit. Trigger: relation to IRIS? -> can be combined as high-quality prior regions.",
            "[20s] Close with four validated contributions and next steps. Trigger: roadmap? -> tighter envelopes, denser obstacles, TAMP integration.",
            "[10s] Thank you and Q&A. Trigger: key takeaway? -> reusable certified free-space maps are the scaling lever.",
        ]
        out_name = "SafeBoxForest_Presentation_v11_en.pptx"
    else:
        notes = [
            "[20s] 主讲：SafeBoxForest 的核心是把一次性搜索升级为可复用体积地图。触发词：贡献？→ 回答‘证书 + 复用 + 增量’。",
            "[20s] 主讲：三大痛点是不可复用、逐点碰撞、动态重建。触发词：为何不用 RRT？→ 回答‘跨查询吞吐劣势’。",
            "[20s] 主讲：三层解耦（证书/森林/规划）保证可扩展。触发词：工程落地？→ 回答‘可独立替换并并行优化’。",
            "[25s] 主讲：流程四步 + 一处关键修复：geodesic + wrap-aware。触发词：路径偏置？→ 回答‘±π 边界代价失真已修复’。",
            "[20s] 主讲：区间 FK 提供保守安全证书，不漏碰撞。触发词：是否过保守？→ 回答‘可能假阳性，但安全优先’。",
            "[20s] 主讲：KD-Tree + promote + cache 是高维效率来源。触发词：为什么快？→ 回答‘复用子问题，减少 FK 重算’。",
            "[20s] 主讲：动态障碍走 invalidate-remove-regrow，仅局部修复。触发词：动态场景能力？→ 回答‘避免全量重建’。",
            "[20s] 主讲：结果只看成功率、速度、适应性三类指标。触发词：路径质量？→ 回答‘geodesic edge + wrap waypoint’。",
            "[20s] 主讲：SBF 与 IRIS/GCS、RRT 是互补定位。触发词：与 IRIS 关系？→ 回答‘可作为高质量先验互补’。",
            "[20s] 主讲：四条贡献收口，并给未来路线。触发词：下一步？→ 回答‘更紧包络、稠密障碍、TAMP’。",
            "[10s] 主讲：谢谢，进入问答。触发词：takeaway？→ 回答‘可复用安全体积地图是规模化关键’。",
        ]
        out_name = "SafeBoxForest_Presentation_v11_zh.pptx"
    _set_notes(prs.slides, notes)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), out_name)
    prs.save(out)
    print(f"Saved: {out}")
    print(f"  {len(prs.slides)} slides (~4 min)")
    return out


if __name__ == "__main__":
    out_zh = build("zh")
    out_en = build("en")
    print(f"Saved bilingual deck pair:\n  - {out_zh}\n  - {out_en}")
